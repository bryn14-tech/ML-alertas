"""Loader de reportes situacionales semanales/mensuales de Antamina (Word/PowerPoint).

Extrae texto de src/data/2025/ y src/data/2026/ (reportes "VP Antamina"),
busca menciones de los 18 distritos de interés (vía la jerarquía de
ancash_datos.py) y etiqueta cada párrafo con: si menciona a Antamina
explícitamente, y si tiene tono de tensión (cuestionan/denuncian/etc.) o de
compromiso/mesa de diálogo.

Es texto narrativo escrito por el equipo de relaciones de Antamina — la
única fuente del proyecto con "compromiso de las partes" explícito (mesas de
diálogo, acuerdos, cronogramas), marcado como "PENDIENTE DE FUENTE" en el
dashboard desde el inicio del proyecto.
"""
import re
import sys
from pathlib import Path

import docx
import pandas as pd
import pptx

from src.scoring.ancash_datos import JERARQUIA, _norm

CARPETAS = [Path("src/data") / "2025", Path("src/data") / "2026"]
INTERIM_DIR = Path("data/interim")
SALIDA = INTERIM_DIR / "reportes_antamina_parrafos.parquet"

PAT_FECHA = re.compile(r"(\d{2})\.(\d{2})\.(\d{2})")
PALABRAS_NEGATIVO = ["cuestion", "denuncia", "protest", "exig", "critic", "alert",
                      "reclam", "rechaz", "bloque", "paraliz", "amenaz"]
PALABRAS_COMPROMISO = ["mesa de dialogo", "compromiso", "acuerdo", "coordina",
                        "avance", "socializ", "cronograma"]


def _extraer_docx(ruta: Path) -> str:
    doc = docx.Document(ruta)
    lineas = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for t in doc.tables:
        for row in t.rows:
            celda = " | ".join(c.text.strip() for c in row.cells)
            if celda.strip(" |"):
                lineas.append(celda)
    return "\n".join(lineas)


def _extraer_pptx(ruta: Path) -> str:
    prs = pptx.Presentation(ruta)
    lineas = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lineas.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    celda = " | ".join(c.text.strip() for c in row.cells)
                    if celda.strip(" |"):
                        lineas.append(celda)
    return "\n".join(lineas)


def _fecha_de_nombre(nombre: str) -> pd.Timestamp | None:
    m = PAT_FECHA.search(nombre)
    if not m:
        return None
    dd, mm, yy = m.groups()
    try:
        return pd.Timestamp(year=2000 + int(yy), month=int(mm), day=int(dd))
    except ValueError:
        return None


def parsear_reportes() -> pd.DataFrame:
    dist2ugt = {_norm(d): ugt for ugt, _, d in JERARQUIA}
    distritos_ordenados = sorted(dist2ugt.keys(), key=len, reverse=True)

    filas = []
    for carpeta in CARPETAS:
        if not carpeta.exists():
            continue
        for archivo in sorted(carpeta.rglob("*")):
            if archivo.suffix.lower() not in (".docx", ".pptx"):
                continue
            fecha = _fecha_de_nombre(archivo.stem)
            if fecha is None:
                continue
            try:
                texto = _extraer_docx(archivo) if archivo.suffix.lower() == ".docx" else _extraer_pptx(archivo)
            except Exception:
                continue

            for parrafo in texto.split("\n"):
                if len(parrafo.strip()) <= 25:
                    continue
                pn = _norm(parrafo)
                ugts_match = {dist2ugt[dn] for dn in distritos_ordenados
                              if re.search(r"\b" + re.escape(dn) + r"\b", pn)}
                if not ugts_match:
                    continue
                menciona_antamina = "antamina" in pn
                tono_negativo = any(k in pn for k in PALABRAS_NEGATIVO)
                tono_compromiso = any(k in pn for k in PALABRAS_COMPROMISO)
                for ugt in ugts_match:
                    filas.append({
                        "fecha_reporte": fecha, "ugt": ugt,
                        "menciona_antamina": menciona_antamina,
                        "tono_negativo": tono_negativo,
                        "tono_compromiso": tono_compromiso,
                    })
    return pd.DataFrame(filas)


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    df = parsear_reportes()
    INTERIM_DIR.mkdir(parents=True, exist_ok=True)
    df.to_parquet(SALIDA, index=False)
    print(f"✓ Reportes Antamina procesados: {len(df)} párrafos con distrito identificado")
    if not df.empty:
        print(f"  Rango: {df['fecha_reporte'].min().date()} → {df['fecha_reporte'].max().date()}")
        print(df.groupby("ugt")[["menciona_antamina", "tono_negativo", "tono_compromiso"]].sum())
    print(f"  Guardado en {SALIDA}")


if __name__ == "__main__":
    main()
